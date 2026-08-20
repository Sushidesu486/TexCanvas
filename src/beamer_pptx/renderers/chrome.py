from __future__ import annotations

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from ..geometry import navigation_widths, page_label
from ..model import Slide, SlideKind
from .common import RenderContext, add_box, add_text, color


def _draw_background(ctx: RenderContext) -> None:
    if ctx.draw_code_background:
        add_box(ctx, "DSH_BACKGROUND", 0, 0, ctx.slide_width, ctx.slide_height, fill=ctx.theme.background)


def _draw_navigation(ctx: RenderContext) -> None:
    count = len(ctx.deck.sections)
    gap = 0.04
    margin = 0.35
    widths = navigation_widths(ctx.slide_width, count, margin=margin, gap=gap)
    font_size = 11 if count <= 5 else 9 if count <= 7 else 7.5
    x = margin
    for index, (section, width) in enumerate(zip(ctx.deck.sections, widths, strict=True)):
        active = index == ctx.section_index
        shape = ctx.slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.CHEVRON,
            Inches(x),
            Inches(0.14),
            Inches(width),
            Inches(0.43),
        )
        shape.name = f"DSH_NAV_{section.id}"
        shape.fill.solid()
        shape.fill.fore_color.rgb = color(ctx.theme.primary if active else ctx.theme.pale)
        shape.line.fill.background()
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = False
        frame.margin_left = frame.margin_right = Inches(0.06)
        frame.margin_top = frame.margin_bottom = Inches(0.01)
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = frame.paragraphs[0]
        paragraph.text = section.short_title
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.runs[0]
        run.font.name = ctx.theme.chinese_font
        run.font.size = Pt(font_size)
        run.font.bold = active
        run.font.color.rgb = color(ctx.theme.white if active else ctx.theme.primary)
        x += width + gap


def _draw_title(ctx: RenderContext, slide: Slide) -> None:
    if slide.kind is SlideKind.SECTION_DIVIDER:
        return
    add_text(
        ctx,
        "DSH_TITLE",
        slide.title,
        0.68,
        0.72,
        ctx.slide_width - 1.36,
        0.55,
        size=24,
        text_color=ctx.theme.primary,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )
    add_box(ctx, "DSH_TITLE_RULE", 0.68, 1.29, ctx.slide_width - 1.36, 0.025, fill=ctx.theme.pale)


def _draw_footer(ctx: RenderContext) -> None:
    y = ctx.slide_height - 0.42
    add_box(ctx, "DSH_FOOTER_RULE", 0.45, y - 0.09, ctx.slide_width - 0.9, 0.025, fill=ctx.theme.secondary)
    left = ctx.deck.metadata.short_title or ctx.deck.metadata.title
    middle = " / ".join(part for part in (ctx.deck.metadata.author, ctx.deck.metadata.institute) if part)
    add_text(ctx, "DSH_FOOTER_LEFT", left, 0.48, y, 4.5, 0.24, size=9, text_color=ctx.theme.muted, margin=0)
    add_text(
        ctx,
        "DSH_FOOTER_MIDDLE",
        middle,
        4.75,
        y,
        ctx.slide_width - 9.5,
        0.24,
        size=9,
        text_color=ctx.theme.muted,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    add_text(
        ctx,
        "DSH_FOOTER_PAGE",
        page_label(ctx.page_number, ctx.total_pages),
        ctx.slide_width - 2.0,
        y,
        1.5,
        0.24,
        size=9,
        text_color=ctx.theme.primary,
        bold=True,
        align=PP_ALIGN.RIGHT,
        margin=0,
    )


def render_chrome(ctx: RenderContext, slide: Slide) -> None:
    _draw_background(ctx)
    _draw_navigation(ctx)
    _draw_title(ctx, slide)
    _draw_footer(ctx)

