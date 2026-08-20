from __future__ import annotations

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from ..assets import inspect_image, resolve_asset
from ..errors import AssetError
from ..geometry import Box, contain, cover
from ..model import Slide
from .common import RenderContext, add_box, add_text


def _placeholder(ctx: RenderContext, message: str, box: Box) -> None:
    add_box(ctx, "DSH_IMAGE_PLACEHOLDER", box.x, box.y, box.width, box.height, fill="EEF2F6", line=ctx.theme.muted)
    add_text(
        ctx,
        "DSH_IMAGE_PLACEHOLDER_TEXT",
        f"Image unavailable\n{message}",
        box.x + 0.3,
        box.y + box.height / 2 - 0.45,
        box.width - 0.6,
        0.9,
        size=14,
        text_color=ctx.theme.muted,
        align=PP_ALIGN.CENTER,
    )


def render_image(ctx: RenderContext, slide: Slide) -> None:
    assert slide.image is not None
    caption_height = 0.45 if slide.caption else 0.1
    target = Box(0.82, 1.52, ctx.slide_width - 1.64, ctx.slide_height - 2.24 - caption_height)
    location = f"sections[{ctx.section_index}].slides.image.path"
    path = resolve_asset(slide.image.path, ctx.asset_root)
    try:
        info = inspect_image(path, location)
    except AssetError as exc:
        if ctx.strict:
            raise
        _placeholder(ctx, str(exc), target)
    else:
        placement = contain(info.width, info.height, target) if slide.image.fit == "contain" else cover(info.width, info.height, target)
        picture = ctx.slide.shapes.add_picture(
            str(info.path),
            Inches(placement.x),
            Inches(placement.y),
            width=Inches(placement.width),
            height=Inches(placement.height),
        )
        picture.name = "DSH_IMAGE"
        picture.crop_left = placement.crop_left
        picture.crop_right = placement.crop_right
        picture.crop_top = placement.crop_top
        picture.crop_bottom = placement.crop_bottom
    if slide.caption:
        add_text(
            ctx,
            "DSH_CAPTION",
            slide.caption,
            0.84,
            ctx.slide_height - 0.98,
            ctx.slide_width - 1.68,
            0.34,
            size=10,
            text_color=ctx.theme.muted,
            align=PP_ALIGN.CENTER,
            margin=0,
        )

