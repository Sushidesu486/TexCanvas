from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from .errors import RenderError
from .model import Deck, SlideKind
from .renderers import (
    render_block,
    render_code,
    render_conclusion,
    render_content,
    render_equation,
    render_figure,
    render_image,
    render_references,
    render_section_divider,
    render_table,
    render_title,
    render_two_columns,
)
from .renderers.chrome import render_chrome
from .renderers.common import EMU_PER_INCH, RenderContext
from .theme import Theme


RENDERERS = {
    SlideKind.TITLE: render_title,
    SlideKind.SECTION_DIVIDER: render_section_divider,
    SlideKind.CONTENT: render_content,
    SlideKind.TWO_COLUMNS: render_two_columns,
    SlideKind.IMAGE: render_image,
    SlideKind.CODE: render_code,
    SlideKind.TABLE: render_table,
    SlideKind.EQUATION: render_equation,
    SlideKind.BLOCK: render_block,
    SlideKind.CONCLUSION: render_conclusion,
    SlideKind.REFERENCES: render_references,
    SlideKind.FIGURE: render_figure,
}


def _blank_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        if layout.name.casefold() in {"blank", "空白"}:
            return layout
    return min(prs.slide_layouts, key=lambda layout: len(layout.placeholders))


def render_deck(
    prs: Presentation,
    deck: Deck,
    *,
    theme: Theme,
    asset_root: Path,
    strict: bool,
    warnings: list[str],
    draw_code_background: bool,
) -> None:
    layout = _blank_layout(prs)
    page_number = 0
    width = prs.slide_width / EMU_PER_INCH
    height = prs.slide_height / EMU_PER_INCH
    for section_index, section in enumerate(deck.sections):
        for slide_model in section.slides:
            page_number += 1
            pptx_slide = prs.slides.add_slide(layout)
            ctx = RenderContext(
                slide=pptx_slide,
                deck=deck,
                section=section,
                section_index=section_index,
                page_number=page_number,
                total_pages=deck.slide_count,
                slide_width=width,
                slide_height=height,
                theme=theme,
                asset_root=asset_root,
                strict=strict,
                warnings=warnings,
                draw_code_background=draw_code_background,
            )
            render_chrome(ctx, slide_model)
            try:
                RENDERERS[slide_model.kind](ctx, slide_model)
            except Exception as exc:
                if isinstance(exc, (RenderError,)):
                    raise
                raise RenderError(f"slide {page_number} ({slide_model.kind.value}): rendering failed: {exc}") from exc

