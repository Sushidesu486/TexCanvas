from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.slide import Slide as PptxSlide
from pptx.util import Inches, Pt

from ..model import Deck, Section
from ..theme import Theme


EMU_PER_INCH = 914400


@dataclass
class RenderContext:
    slide: PptxSlide
    deck: Deck
    section: Section
    section_index: int
    page_number: int
    total_pages: int
    slide_width: float
    slide_height: float
    theme: Theme
    asset_root: Path
    strict: bool
    warnings: list[str]
    draw_code_background: bool


def color(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_run_font(run, *, latin: str, ea: str) -> None:
    """Authoritatively set both the Latin and East-Asian typefaces on a run.

    python-pptx's ``font.name`` only writes ``a:latin``; WPS/PowerPoint pick the
    East-Asian font (used for CJK characters) from ``a:ea`` and the complex-script
    font from ``a:cs``, both of which would otherwise be inherited from the theme.
    Writing all three keeps CJK glyphs in 苹方-简 and Latin/digits/symbols in Helvetica
    regardless of the deck template's theme fonts.
    """
    run.font.name = latin
    rPr = run.font._rPr
    for tag, typeface in (("a:ea", ea), ("a:cs", ea)):
        existing = rPr.find(qn(tag))
        if existing is None:
            existing = rPr.makeelement(qn(tag), {})
            rPr.append(existing)
        existing.set("typeface", typeface)



def add_box(
    ctx: RenderContext,
    name: str,
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    fill: str,
    line: str | None = None,
    radius: bool = False,
):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = ctx.slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(width), Inches(height))
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = color(line)
    return shape


def add_text(
    ctx: RenderContext,
    name: str,
    text: str,
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    size: float = 20,
    text_color: str | None = None,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.04,
    font_name: str | None = None,
):
    shape = ctx.slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    shape.name = name
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.1
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    latin = font_name or ctx.theme.chinese_font
    set_run_font(run, latin=latin, ea=ctx.theme.chinese_font)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color(text_color or ctx.theme.text)
    return shape


def add_rich_text(
    ctx: RenderContext,
    name: str,
    paragraphs: list[tuple[str, float, bool, str, float]],
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    margin: float = 0.08,
):
    """Add paragraphs as (text, font size, bold, color, space_after_pt)."""
    shape = ctx.slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    shape.name = name
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    for index, (value, size, bold, text_color, space_after) in enumerate(paragraphs):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = value
        paragraph.space_after = Pt(space_after)
        paragraph.line_spacing = 1.12
        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        set_run_font(run, latin=ctx.theme.chinese_font, ea=ctx.theme.chinese_font)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color(text_color)
    return shape


def content_height(ctx: RenderContext) -> float:
    return max(0.5, ctx.slide_height - 1.48 - 0.58)
