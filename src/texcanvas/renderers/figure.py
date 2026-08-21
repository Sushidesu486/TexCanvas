from __future__ import annotations

import tempfile
from pathlib import Path

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from ..errors import RenderError, TexCanvasError
from ..geometry import Box, contain, cover
from ..graphics import FigureSpec as GSpec, _build_figure
from ..model import Slide
from .common import RenderContext, add_text
from .image import _placeholder


def _compile_inline_figure(spec, slide_id: str, output_dir: Path) -> Path:
    """Compile an inline deck figure spec to a PNG via the graphics pipeline.

    The graphics module builds from a file-backed spec; for an inline spec we
    write the source text to a temp file and let the existing pipeline compile
    it.  Only PNG is produced (the PPTX needs a raster).
    """
    with tempfile.TemporaryDirectory(prefix=f"texcanvas-fig-{slide_id}-") as temp_dir:
        temp = Path(temp_dir)
        src_path = temp / f"{slide_id}.tikz"
        src_path.write_text(spec.source, encoding="utf-8")
        gspec = GSpec(
            id=slide_id,
            engine=spec.engine,
            source=src_path,
            outputs=("png",),
            preamble=spec.preamble,
            compiler=spec.compiler,
            dpi=spec.dpi,
        )
        generated = _build_figure(gspec, output_dir)
        png = generated[0]
        # The temp dir holding src_path closes on exit; _build_figure wrote the
        # PNG to output_dir (outside the temp dir), so it survives.
        return png


def render_figure(ctx: RenderContext, slide: Slide) -> None:
    assert slide.figure is not None
    spec = slide.figure
    caption_height = 0.45 if slide.caption else 0.1
    target = Box(0.82, 1.52, ctx.slide_width - 1.64, ctx.slide_height - 2.24 - caption_height)
    slide_id = f"slide{ctx.page_number}"

    # Figures compile into a sibling of the deck output so repeated builds reuse them.
    fig_dir = Path(str(ctx.asset_root)) / "_texcanvas_figures"
    fig_dir.mkdir(parents=True, exist_ok=True)

    try:
        png_path = _compile_inline_figure(spec, slide_id, fig_dir)
    except TexCanvasError:
        if ctx.strict:
            raise
        _placeholder(ctx, "figure compile failed", target)
        png_path = None
    except Exception as exc:
        if ctx.strict:
            raise RenderError(f"slide {ctx.page_number} (figure): compile failed: {exc}") from exc
        _placeholder(ctx, str(exc), target)
        png_path = None

    if png_path is not None:
        from ..assets import inspect_image

        location = f"sections[{ctx.section_index}].slides.figure"
        try:
            info = inspect_image(png_path, location)
        except TexCanvasError:
            if ctx.strict:
                raise
            _placeholder(ctx, "compiled figure unreadable", target)
        else:
            placement = contain(info.width, info.height, target) if spec.fit == "contain" else cover(info.width, info.height, target)
            picture = ctx.slide.shapes.add_picture(
                str(info.path),
                Inches(placement.x),
                Inches(placement.y),
                width=Inches(placement.width),
                height=Inches(placement.height),
            )
            picture.name = "DSH_FIGURE"
            picture.crop_left = placement.crop_left
            picture.crop_right = placement.crop_right
            picture.crop_top = placement.crop_top
            picture.crop_bottom = placement.crop_bottom

    if slide.caption:
        add_text(
            ctx,
            "DSH_CAPTION",
            slide.caption,
            0.84,
            ctx.slide_height - 0.98,
            ctx.slide_width - 1.68,
            0.34,
            size=10,
            text_color=ctx.theme.muted,
            align=PP_ALIGN.CENTER,
            margin=0,
        )
