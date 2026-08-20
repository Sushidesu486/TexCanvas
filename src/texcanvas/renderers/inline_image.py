from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from ..assets import inspect_image, resolve_asset
from ..errors import AssetError
from ..geometry import Box, ImagePlacement, contain
from ..model import InlineImage, Slide
from .common import RenderContext


def _fit_inline(image_width: int, image_height: int, target_width: float) -> tuple[float, float]:
    """Scale an image to ``target_width`` inches, preserving aspect ratio."""
    if image_width <= 0:
        return target_width, target_width
    scale = target_width / image_width
    return target_width, image_height * scale


def render_inline_image(ctx: RenderContext, slide: Slide, content_box: Box) -> Box:
    """Place the slide's inline image and return the text-safe box.

    The image is anchored to one side of ``content_box``; the returned box is
    the remaining area the caller should keep text within so it does not
    overlap the image. On any asset error under ``--no-strict`` the image is
    skipped and the full ``content_box`` is returned.
    """
    assert slide.inline_image is not None
    spec = slide.inline_image
    location = f"sections[{ctx.section_index}].slides.inline_image.path"
    path = resolve_asset(spec.path, ctx.asset_root)
    try:
        info = inspect_image(path, location)
    except AssetError:
        if ctx.strict:
            raise
        return content_box

    img_w_in, img_h_in = _fit_inline(info.width, info.height, spec.width)
    # Clamp the image height so it never exceeds the content box height.
    if img_h_in > content_box.height:
        img_h_in = content_box.height
        img_w_in = img_h_in * (info.width / info.height)
    top_pad = max(0.0, (content_box.height - img_h_in) / 2)

    if spec.align == "right":
        img_x = content_box.x + content_box.width - img_w_in
        text_box = Box(content_box.x, content_box.y, content_box.width - img_w_in - 0.25, content_box.height)
    else:
        img_x = content_box.x
        text_box = Box(content_box.x + img_w_in + 0.25, content_box.y, content_box.width - img_w_in - 0.25, content_box.height)

    picture = ctx.slide.shapes.add_picture(
        str(info.path),
        Inches(img_x),
        Inches(content_box.y + top_pad),
        width=Inches(img_w_in),
        height=Inches(img_h_in),
    )
    picture.name = "DSH_INLINE_IMAGE"
    return text_box
