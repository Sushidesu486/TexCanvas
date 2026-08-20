from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn

from ..model import Slide
from .common import RenderContext, add_text, set_run_font


def _set_cell_margins(cell, *, left: float = 0.08, right: float = 0.08, top: float = 0.03, bottom: float = 0.03) -> None:
    """Set uniform cell margins (inches) via the tcPr XML element."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for attr, value in (("marL", int(left * 914400)), ("marR", int(right * 914400)),
                        ("marT", int(top * 914400)), ("marB", int(bottom * 914400))):
        tcPr.set(qn(f"a:{attr}"), str(value))


def _style_cell(cell, *, fill: str, text_color: str, font_name: str, ea: str, size: float, bold: bool, align: PP_ALIGN) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor.from_string(fill)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Inches(0.1)
    cell.margin_top = cell.margin_bottom = Inches(0.04)
    frame = cell.text_frame
    frame.clear()
    frame.word_wrap = True
    for paragraph in frame.paragraphs:
        paragraph.alignment = align
        paragraph.line_spacing = 1.1
        paragraph.space_after = Pt(0)
    run = frame.paragraphs[0].add_run()
    set_run_font(run, latin=font_name, ea=ea)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(text_color)


def _set_cell_text(cell, text: str) -> None:
    cell.text_frame.paragraphs[0].runs[0].text = text if text is not None else ""


def render_table(ctx: RenderContext, slide: Slide) -> None:
    assert slide.table is not None
    table_spec = slide.table
    header = table_spec.header
    rows = table_spec.rows
    body_rows = list(rows)

    caption_height = 0.42 if slide.caption or table_spec.caption else 0.0
    table_top = 1.55
    available_height = ctx.slide_height - 1.48 - 0.6 - table_top - caption_height

    has_header = bool(header)
    total_rows = (len(body_rows) + (1 if has_header else 0))
    total_rows = max(total_rows, 1)

    # Distribute height; header slightly taller than body rows.
    header_factor = 1.25 if has_header else 0.0
    body_units = max(1, len(body_rows)) * 1.0 + header_factor
    body_row_height = min(0.6, max(0.22, available_height / body_units)) if body_rows else 0.3
    header_row_height = body_row_height * header_factor if has_header else 0.0
    table_height = header_row_height + body_row_height * len(body_rows)

    # Column count: prefer header, else max row length.
    col_count = max(len(header) if has_header else 0, *(len(row) for row in body_rows), 1)
    table_width = ctx.slide_width - 1.64
    table_left = 0.82

    nrows = (1 if has_header else 0) + len(body_rows)
    nrows = max(nrows, 1)
    shape = ctx.slide.shapes.add_table(nrows, col_count, Inches(table_left), Inches(table_top), Inches(table_width), Inches(table_height))
    table = shape.table
    shape.name = "DSH_TABLE"

    # Equal column widths.
    col_width = table_width / col_count
    for col_index in range(col_count):
        table.columns[col_index].width = Inches(col_width)

    # Row heights.
    first_body = 0 if not has_header else 1
    if has_header:
        table.rows[0].height = Inches(header_row_height)
    for body_index in range(len(body_rows)):
        table.rows[first_body + body_index].height = Inches(body_row_height)

    # Disable the built-in banded style so our explicit fills win in WPS.
    _disable_table_style(table)

    if has_header:
        for col_index, label in enumerate(header):
            cell = table.cell(0, col_index)
            _style_cell(
                cell,
                fill=ctx.theme.table_header_fill,
                text_color=ctx.theme.table_header_text,
                font_name=ctx.theme.chinese_font,
                ea=ctx.theme.chinese_font,
                size=14,
                bold=True,
                align=PP_ALIGN.CENTER,
            )
            _set_cell_text(cell, label)

    body_size = 13 if len(body_rows) <= 8 else 12 if len(body_rows) <= 10 else 11
    for row_index, row in enumerate(body_rows):
        grid_row = first_body + row_index
        alt = row_index % 2 == 1
        fill = ctx.theme.table_row_alt if alt else ctx.theme.white
        for col_index in range(col_count):
            text = row[col_index] if col_index < len(row) else ""
            cell = table.cell(grid_row, col_index)
            _style_cell(
                cell,
                fill=fill,
                text_color=ctx.theme.text,
                font_name=ctx.theme.chinese_font,
                ea=ctx.theme.chinese_font,
                size=body_size,
                bold=False,
                align=PP_ALIGN.LEFT,
            )
            _set_cell_text(cell, text)

    # Explicit per-cell borders give a clean grid in WPS without relying on a table style.
    _apply_cell_borders(table, ctx.theme.table_grid, has_header, len(body_rows), col_count)

    caption = slide.caption or table_spec.caption
    if caption:
        add_text(
            ctx,
            "DSH_CAPTION",
            caption,
            0.84,
            ctx.slide_height - 0.98,
            ctx.slide_width - 1.68,
            0.34,
            size=10,
            text_color=ctx.theme.muted,
            align=PP_ALIGN.CENTER,
            margin=0,
        )


def _disable_table_style(table) -> None:
    """Remove the first-row banding so explicit per-cell fills survive in WPS."""
    tbl = table._tbl
    tblPr = tbl.find(qn("a:tblPr"))
    if tblPr is None:
        tblPr = tbl.makeelement(qn("a:tblPr"), {})
        tbl.insert(0, tblPr)
    tblPr.set("firstRow", "0")
    tblPr.set("bandRow", "0")


def _apply_cell_borders(table, color: str, has_header: bool, body_count: int, col_count: int) -> None:
    last_row = (1 if has_header else 0) + body_count - 1
    for row_index in range(last_row + 1):
        for col_index in range(col_count):
            _border_cell(table.cell(row_index, col_index), color)


def _border_cell(cell, color: str) -> None:
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # Remove any existing line definitions.
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        existing = tcPr.find(qn(tag))
        if existing is not None:
            tcPr.remove(existing)
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        ln = tcPr.makeelement(qn(tag), {"w": "9525", "cap": "flat"})
        fill = ln.makeelement(qn("a:solidFill"), {})
        srgb = fill.makeelement(qn("a:srgbClr"), {"val": color})
        fill.append(srgb)
        ln.append(fill)
        tcPr.append(ln)
