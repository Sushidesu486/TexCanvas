"""Extract and apply human edits from generated PPTX files.

The sync format is intentionally a presentation override layer rather than a
second source of semantic slide content.  YAML remains authoritative for the
deck model; ``overrides.yml`` carries geometry, text edits, slide order, and
new pictures discovered in a human-edited PPTX.
"""

from __future__ import annotations

import posixpath
import re
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any

import yaml
from lxml import etree
from pptx import Presentation
from pptx.util import Inches

from .errors import InputError, ValidationError

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
PKG_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"
EMU_PER_INCH = 914400

_SLIDE_REL_TYPE = "/slide"


@dataclass(frozen=True)
class SyncReport:
    output: Path
    slide_count: int
    shape_count: int
    extracted_image_count: int


def pull(input: str | Path, output: str | Path, asset_dir: str | Path | None = None) -> SyncReport:
    """Extract a PPTX's current presentation state into an overrides YAML file."""
    input_path = Path(input).expanduser().resolve()
    output_path = Path(output).expanduser().resolve()
    if not input_path.is_file():
        raise InputError(f"sync input: file not found: {input_path}")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    assets_path = Path(asset_dir).expanduser().resolve() if asset_dir is not None else output_path.parent / f"{output_path.stem}-assets"
    assets_path.mkdir(parents=True, exist_ok=True)

    try:
        with zipfile.ZipFile(input_path) as archive:
            entries = {name: archive.read(name) for name in archive.namelist()}
            slide_parts = _ordered_slide_parts(entries)
            slides: list[dict[str, Any]] = []
            shape_count = 0
            image_count = 0
            for order, slide_path in enumerate(slide_parts):
                root = etree.fromstring(entries[slide_path])
                slide_id = _slide_id(root, order)
                relationships = _relationships(entries, slide_path)
                shape_map: dict[str, Any] = {}
                for shape_index, shape in enumerate(_iter_shapes(root)):
                    snapshot, image_extracted = _shape_snapshot(
                        shape,
                        slide_id=slide_id,
                        shape_index=shape_index,
                        relationships=relationships,
                        entries=entries,
                        assets_path=assets_path,
                        output_parent=output_path.parent,
                    )
                    key = snapshot.pop("name")
                    if key in shape_map:
                        key = f"{key}#{shape_index + 1}"
                    shape_map[key] = snapshot
                    shape_count += 1
                    image_count += image_extracted
                slides.append({"id": slide_id, "order": order, "shapes": shape_map})
    except (KeyError, OSError, zipfile.BadZipFile, etree.XMLSyntaxError) as exc:
        raise InputError(f"sync input: cannot inspect {input_path}: {exc}") from exc

    payload = {"version": 1, "slides": slides}
    output_path.write_text(yaml.safe_dump(payload, allow_unicode=True, sort_keys=False), encoding="utf-8")
    return SyncReport(output_path, len(slides), shape_count, image_count)


def apply_overrides(prs: Presentation, path: str | Path) -> None:
    """Apply an overrides YAML file to an in-memory generated presentation."""
    source = Path(path).expanduser().resolve()
    try:
        payload = yaml.safe_load(source.read_text(encoding="utf-8"))
    except (OSError, yaml.YAMLError) as exc:
        raise InputError(f"overrides: cannot read {source}: {exc}") from exc
    if not isinstance(payload, dict) or payload.get("version") != 1:
        raise ValidationError("overrides.version: expected 1")
    slides_data = payload.get("slides")
    if not isinstance(slides_data, list):
        raise ValidationError("overrides.slides: expected a list")

    current_slides = list(prs.slides)
    slide_by_id = {_slide_identity(slide, index): slide for index, slide in enumerate(current_slides)}
    _reorder_slides(prs, slide_by_id, slides_data)

    for slide_data in slides_data:
        if not isinstance(slide_data, dict):
            continue
        slide = slide_by_id.get(str(slide_data.get("id", "")))
        if slide is None:
            continue
        shapes = slide_data.get("shapes", {})
        if not isinstance(shapes, dict):
            continue
        existing = {shape.name: shape for shape in slide.shapes}
        for name, spec in shapes.items():
            if not isinstance(spec, dict):
                continue
            shape = existing.get(name)
            if shape is None and spec.get("kind") == "picture":
                asset = spec.get("asset")
                if not isinstance(asset, str):
                    raise ValidationError(f"overrides.slides[{slide_data.get('id')}].shapes[{name}].asset: required")
                asset_path = (source.parent / asset).resolve()
                if not asset_path.is_file():
                    raise InputError(f"overrides asset: file not found: {asset_path}")
                frame = _frame_values(spec.get("frame"))
                frame_emu = spec.get("frame_emu")
                if isinstance(frame_emu, dict) and all(key in frame_emu for key in ("x", "y", "width", "height")):
                    frame = {key: int(frame_emu[key]) / EMU_PER_INCH for key in ("x", "y", "width", "height")}
                shape = slide.shapes.add_picture(
                    str(asset_path), Inches(frame["x"]), Inches(frame["y"]), Inches(frame["width"]), Inches(frame["height"])
                )
                shape.name = name
                existing[name] = shape
            if shape is None:
                continue
            _apply_frame(shape, spec)
            if "text" in spec and hasattr(shape, "text_frame"):
                _apply_text(shape, spec)


def _ordered_slide_parts(entries: dict[str, bytes]) -> list[str]:
    presentation = etree.fromstring(entries["ppt/presentation.xml"])
    relationships = etree.fromstring(entries["ppt/_rels/presentation.xml.rels"])
    relation_targets = {
        rel.get("Id"): posixpath.normpath(posixpath.join("ppt", rel.get("Target", "")))
        for rel in relationships
        if rel.get("Type", "").endswith(_SLIDE_REL_TYPE)
    }
    result: list[str] = []
    slide_ids = presentation.find("./{%s}sldIdLst" % P_NS)
    if slide_ids is None:
        return result
    for slide_id in slide_ids:
        target = relation_targets.get(slide_id.get("{%s}id" % R_NS))
        if target in entries:
            result.append(target)
    return result


def _relationships(entries: dict[str, bytes], slide_path: str) -> dict[str, str]:
    rels_path = posixpath.join(posixpath.dirname(slide_path), "_rels", posixpath.basename(slide_path) + ".rels")
    if rels_path not in entries:
        return {}
    root = etree.fromstring(entries[rels_path])
    return {
        rel.get("Id"): posixpath.normpath(posixpath.join(posixpath.dirname(slide_path), rel.get("Target", "")))
        for rel in root
        if rel.get("Id") and rel.get("Target")
    }


def _slide_id(root: etree._Element, index: int) -> str:
    c_sld = root.find("./{%s}cSld" % P_NS)
    value = c_sld.get("name", "") if c_sld is not None else ""
    if value.startswith("texcanvas:"):
        return value.removeprefix("texcanvas:")
    return value or f"slide-{index + 1}"


def _slide_identity(slide: Any, index: int) -> str:
    c_sld = slide._element.cSld
    value = c_sld.get("name", "") if c_sld is not None else ""
    if value.startswith("texcanvas:"):
        return value.removeprefix("texcanvas:")
    return value or f"slide-{index + 1}"


def _iter_shapes(root: etree._Element):
    sp_tree = root.find(".//{%s}spTree" % P_NS)
    if sp_tree is None:
        return
    for child in sp_tree:
        if child.tag in {_q(P_NS, "sp"), _q(P_NS, "pic"), _q(P_NS, "graphicFrame")}:
            yield child
        elif child.tag == _q(MC_NS, "AlternateContent"):
            choice = child.find("{%s}Choice" % MC_NS)
            if choice is not None:
                for nested in choice:
                    if nested.tag in {_q(P_NS, "sp"), _q(P_NS, "pic"), _q(P_NS, "graphicFrame")}:
                        yield nested


def _shape_snapshot(
    shape: etree._Element,
    *,
    slide_id: str,
    shape_index: int,
    relationships: dict[str, str],
    entries: dict[str, bytes],
    assets_path: Path,
    output_parent: Path,
) -> tuple[dict[str, Any], int]:
    c_nv_pr = shape.find("./{%s}nvSpPr/{%s}cNvPr" % (P_NS, P_NS))
    if c_nv_pr is None:
        c_nv_pr = shape.find("./{%s}nvPicPr/{%s}cNvPr" % (P_NS, P_NS))
    name = c_nv_pr.get("name") if c_nv_pr is not None else f"shape-{shape_index + 1}"
    kind = {"{%s}sp" % P_NS: "shape", "{%s}pic" % P_NS: "picture", "{%s}graphicFrame" % P_NS: "graphic"}.get(shape.tag, "shape")
    snapshot: dict[str, Any] = {"name": name, "kind": kind, "generated": name.startswith("DSH_")}
    frame = shape.find(".//{%s}xfrm" % A_NS)
    if frame is not None:
        off = frame.find("{%s}off" % A_NS)
        ext = frame.find("{%s}ext" % A_NS)
        if off is not None and ext is not None:
            snapshot["frame"] = {
                "x": round(int(off.get("x", "0")) / EMU_PER_INCH, 4),
                "y": round(int(off.get("y", "0")) / EMU_PER_INCH, 4),
                "width": round(int(ext.get("cx", "0")) / EMU_PER_INCH, 4),
                "height": round(int(ext.get("cy", "0")) / EMU_PER_INCH, 4),
            }
            snapshot["frame_emu"] = {
                "x": int(off.get("x", "0")),
                "y": int(off.get("y", "0")),
                "width": int(ext.get("cx", "0")),
                "height": int(ext.get("cy", "0")),
            }
    tx_body = shape.find(".//{%s}txBody" % P_NS)
    if tx_body is not None:
        paragraphs = []
        for paragraph in tx_body.findall("{%s}p" % A_NS):
            paragraphs.append("".join(node.text or "" for node in paragraph.iter("{%s}t" % A_NS)))
        if any(paragraphs):
            snapshot["paragraphs"] = paragraphs
            snapshot["text"] = "\n".join(paragraphs)

    extracted = 0
    if kind == "picture":
        blip = shape.find(".//{%s}blip" % A_NS)
        rel_id = blip.get("{%s}embed" % R_NS) if blip is not None else None
        image_path = relationships.get(rel_id or "")
        if image_path in entries:
            safe_slide = re.sub(r"[^A-Za-z0-9_.-]+", "-", slide_id)
            safe_name = re.sub(r"[^A-Za-z0-9_.-]+", "-", name)
            suffix = Path(image_path).suffix or ".bin"
            asset_name = f"{safe_slide}-{safe_name}{suffix}"
            assets_path.mkdir(parents=True, exist_ok=True)
            asset_path = assets_path / asset_name
            asset_path.write_bytes(entries[image_path])
            snapshot["asset"] = Path(posixpath.relpath(asset_path, output_parent)).as_posix()
            extracted = 1
    return snapshot, extracted


def _q(namespace: str, tag: str) -> str:
    return "{%s}%s" % (namespace, tag)


def _frame_values(value: Any) -> dict[str, float]:
    if not isinstance(value, dict):
        raise ValidationError("overrides.frame: expected a mapping")
    return {key: float(value[key]) for key in ("x", "y", "width", "height")}


def _apply_frame(shape: Any, value: Any) -> None:
    if value is None:
        return
    frame_value = value.get("frame") if isinstance(value, dict) and "frame" in value else value
    frame = _frame_values(frame_value)
    emu = value.get("frame_emu") if isinstance(value, dict) else None
    if isinstance(emu, dict) and all(key in emu for key in ("x", "y", "width", "height")):
        shape.left = int(emu["x"])
        shape.top = int(emu["y"])
        shape.width = int(emu["width"])
        shape.height = int(emu["height"])
    else:
        shape.left = Inches(frame["x"])
        shape.top = Inches(frame["y"])
        shape.width = Inches(frame["width"])
        shape.height = Inches(frame["height"])


def _apply_text(shape: Any, spec: dict[str, Any]) -> None:
    frame = shape.text_frame
    paragraphs = spec.get("paragraphs")
    if not isinstance(paragraphs, list):
        paragraphs = [str(spec.get("text", ""))]
    current = [paragraph.text for paragraph in frame.paragraphs]
    if current[: len(paragraphs)] == [str(value) for value in paragraphs] and len(current) == len(paragraphs):
        return
    for index, value in enumerate(paragraphs):
        paragraph = frame.paragraphs[index] if index < len(frame.paragraphs) else frame.add_paragraph()
        text = str(value)
        if paragraph.runs:
            paragraph.runs[0].text = text
            for run in list(paragraph.runs[1:]):
                run._r.getparent().remove(run._r)
        else:
            paragraph.text = text
    while len(frame.paragraphs) > len(paragraphs):
        paragraph = frame.paragraphs[-1]
        paragraph._p.getparent().remove(paragraph._p)


def _reorder_slides(prs: Presentation, slide_by_id: dict[str, Any], slides_data: list[Any]) -> None:
    desired = [item for item in sorted(slides_data, key=lambda item: item.get("order", 0)) if isinstance(item, dict)]
    current_sld_ids = prs.slides._sldIdLst
    pairs = list(zip(list(prs.slides), list(current_sld_ids)))
    sld_id_by_identity = {
        _slide_identity(slide, index): sld_id for index, (slide, sld_id) in enumerate(pairs)
    }
    ordered_ids = [sld_id_by_identity[item.get("id", "")] for item in desired if item.get("id", "") in sld_id_by_identity]
    if not ordered_ids:
        return
    for sld_id in list(current_sld_ids):
        current_sld_ids.remove(sld_id)
    for sld_id in ordered_ids:
        current_sld_ids.append(sld_id)
