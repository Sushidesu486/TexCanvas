from pathlib import Path

from pptx import Presentation

from beamer_pptx import build


def test_build_returns_report_and_reopens(deck_files: tuple[Path, Path], tmp_path: Path):
    source, asset_root = deck_files
    output = tmp_path / "result.pptx"
    report = build(source, output, asset_root=asset_root)
    assert report.slide_count == 6
    assert report.section_count == 3
    assert report.warnings == ()
    assert output.is_file() and output.stat().st_size > 0
    assert len(Presentation(output).slides) == 6


def test_no_strict_uses_visible_placeholder(tmp_path: Path):
    source = tmp_path / "missing.yml"
    source.write_text(
        "metadata: {title: Missing}\nsections:\n  - title: S\n    slides:\n      - kind: image\n        title: Image\n        image: {path: absent.png}\n",
        encoding="utf-8",
    )
    output = tmp_path / "missing.pptx"
    report = build(source, output, strict=False)
    assert report.warnings
    names = {shape.name for shape in Presentation(output).slides[0].shapes}
    assert "DSH_IMAGE_PLACEHOLDER" in names

