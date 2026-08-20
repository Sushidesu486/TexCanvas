from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation

from texcanvas import build


def _shape(slide, name):
    return next(shape for shape in slide.shapes if shape.name == name)


def _make_png(path: Path) -> Path:
    Image.new("RGB", (640, 320), "#2A6E3F").save(path)
    return path


def test_inline_image_right_align_narrows_body(tmp_path: Path):
    asset = _make_png(tmp_path / "fig.png")
    source = tmp_path / "deck.yml"
    source.write_text(
        """
metadata: {title: Inline}
sections:
  - id: a
    title: A
    short_title: A
    slides:
      - kind: content
        title: Body
        body: Some body text.
        bullets: [one, two]
        inline_image: {path: fig.png, width: 3.0, align: right}
""".strip(),
        encoding="utf-8",
    )
    output = tmp_path / "out.pptx"
    build(source, output, asset_root=tmp_path)
    prs = Presentation(output)
    slide = prs.slides[0]
    body = _shape(slide, "DSH_BODY")
    picture = _shape(slide, "DSH_INLINE_IMAGE")
    # Image sits on the right edge of the content box.
    content_right = 0.82 + (13.333 - 1.64)
    img_right = (picture.left + picture.width) / 914400
    assert abs(img_right - content_right) < 0.05
    # Body width shrunk to leave room for the image plus a gap.
    assert body.width / 914400 < (13.333 - 1.64) - 3.0


def test_inline_image_left_align_puts_text_on_right(tmp_path: Path):
    _make_png(tmp_path / "fig.png")
    source = tmp_path / "deck.yml"
    source.write_text(
        """
metadata: {title: Inline}
sections:
  - id: a
    title: A
    short_title: A
    slides:
      - kind: content
        title: Body
        body: Text.
        inline_image: {path: fig.png, width: 3.0, align: left}
""".strip(),
        encoding="utf-8",
    )
    output = tmp_path / "out.pptx"
    build(source, output, asset_root=tmp_path)
    prs = Presentation(output)
    slide = prs.slides[0]
    body = _shape(slide, "DSH_BODY")
    picture = _shape(slide, "DSH_INLINE_IMAGE")
    # Image on the left edge of content box; body starts after it.
    assert picture.left / 914400 < 1.0
    assert body.left / 914400 > picture.left / 914400 + picture.width / 914400


def test_inline_image_missing_strict_fails(tmp_path: Path):
    source = tmp_path / "deck.yml"
    source.write_text(
        """
metadata: {title: Inline}
sections:
  - id: a
    title: A
    short_title: A
    slides:
      - kind: content
        title: Body
        body: Text.
        inline_image: {path: absent.png, align: right}
""".strip(),
        encoding="utf-8",
    )
    output = tmp_path / "out.pptx"
    import pytest
    from texcanvas.errors import AssetError

    with pytest.raises(AssetError):
        build(source, output, asset_root=tmp_path, strict=True)


def test_inline_image_missing_no_strict_warns_and_keeps_full_body(tmp_path: Path):
    source = tmp_path / "deck.yml"
    source.write_text(
        """
metadata: {title: Inline}
sections:
  - id: a
    title: A
    short_title: A
    slides:
      - kind: content
        title: Body
        body: Text.
        inline_image: {path: absent.png, align: right}
""".strip(),
        encoding="utf-8",
    )
    output = tmp_path / "out.pptx"
    report = build(source, output, asset_root=tmp_path, strict=False)
    assert report.warnings
    prs = Presentation(output)
    slide = prs.slides[0]
    names = {sh.name for sh in slide.shapes}
    assert "DSH_INLINE_IMAGE" not in names  # image skipped
    # Body keeps the full content width since no image was placed.
    body = _shape(slide, "DSH_BODY")
    assert body.width / 914400 > (13.333 - 1.64) - 0.1


def test_citation_renders_grey_strip(tmp_path: Path):
    source = tmp_path / "deck.yml"
    source.write_text(
        """
metadata: {title: Cite}
sections:
  - id: a
    title: A
    short_title: A
    slides:
      - kind: content
        title: Body
        body: Text.
        citation: "[1] Doe et al., NeurIPS 2024."
""".strip(),
        encoding="utf-8",
    )
    output = tmp_path / "out.pptx"
    build(source, output, asset_root=tmp_path)
    prs = Presentation(output)
    slide = prs.slides[0]
    citation = _shape(slide, "DSH_CITATION")
    assert citation.text_frame.text == "[1] Doe et al., NeurIPS 2024."
    run = citation.text_frame.paragraphs[0].runs[0]
    # Grey (muted) and small.
    assert run.font.size.pt == 9
    assert str(run.font.color.rgb) == "65727E"
    # Sits above the footer rule.
    footer = _shape(slide, "DSH_FOOTER_RULE")
    assert citation.top < footer.top


def test_citation_absent_when_not_set(tmp_path: Path):
    source = tmp_path / "deck.yml"
    source.write_text(
        """
metadata: {title: Cite}
sections:
  - id: a
    title: A
    short_title: A
    slides:
      - kind: content
        title: Body
        body: Text.
""".strip(),
        encoding="utf-8",
    )
    output = tmp_path / "out.pptx"
    build(source, output, asset_root=tmp_path)
    prs = Presentation(output)
    names = {sh.name for sh in prs.slides[0].shapes}
    assert "DSH_CITATION" not in names
