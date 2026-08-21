from __future__ import annotations

import zipfile
from pathlib import Path

from lxml import etree
import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from texcanvas import build
from texcanvas.mathml import normalize_math_namespaces_in_pptx


def _shape(slide, name):
    return next(shape for shape in slide.shapes if shape.name == name)


def _equation_omath(output: Path, slide_number: int):
    math_ns = "http://schemas.openxmlformats.org/officeDocument/2006/math"
    with zipfile.ZipFile(output) as archive:
        root = etree.fromstring(archive.read(f"ppt/slides/slide{slide_number}.xml"))
    omath = root.find(".//{%s}oMath" % math_ns)
    assert omath is not None
    return omath


@pytest.fixture
def primitives_deck(tmp_path: Path) -> tuple[Path, Path]:
    assets = tmp_path / "assets"
    assets.mkdir()
    yaml_path = tmp_path / "deck.yml"
    yaml_path.write_text(
        r"""
metadata: {title: Primitives, short_title: Prim}
aspect: "16:9"
sections:
  - id: code
    title: Code
    short_title: Code
    slides:
      - kind: code
        title: Python
        code: {lang: python, source: "def f(x):\n    return x + 1\n"}
  - id: table
    title: Table
    short_title: Table
    slides:
      - kind: table
        title: Results
        table:
          header: [Method, Score]
          rows: [[Baseline, "0.8"], [Ours, "0.9"]]
          caption: Table 1
  - id: equation
    title: Equation
    short_title: Eq
    slides:
      - kind: equation
        title: Loss
        equation: |
          L = -\frac{1}{N}\sum_{i=1}^{N} y_i
  - id: block
    title: Blocks
    short_title: Blocks
    slides:
      - kind: block
        title: Note
        block: {style: default, title: Note, body: A note, bullets: [one, two]}
      - kind: block
        title: Alert
        block: {style: alert, title: Warning, body: Be careful}
      - kind: block
        title: Example
        block: {style: example, title: Demo, body: An example}
""".strip(),
        encoding="utf-8",
    )
    return yaml_path, tmp_path


def test_build_with_new_kinds_and_reopens(primitives_deck, tmp_path: Path):
    source, asset_root = primitives_deck
    output = tmp_path / "out.pptx"
    report = build(source, output, asset_root=asset_root)
    assert report.slide_count == 6
    assert report.section_count == 4
    assert report.warnings == ()
    assert output.is_file() and output.stat().st_size > 0
    assert len(Presentation(output).slides) == 6


def test_code_renders_with_panel_and_header_runs(primitives_deck, tmp_path: Path):
    source, asset_root = primitives_deck
    output = tmp_path / "code.pptx"
    build(source, output, asset_root=asset_root)
    prs = Presentation(output)
    code_slide = prs.slides[0]
    panel = _shape(code_slide, "DSH_CODE_PANEL")
    assert panel is not None
    body = _shape(code_slide, "DSH_CODE_BODY")
    # Line 1: "def f(x):" -> def keyword; Line 2: "    return x + 1" -> return keyword.
    first_line_runs = body.text_frame.paragraphs[0].runs
    second_line_runs = body.text_frame.paragraphs[1].runs
    assert any(run.text == "def" for run in first_line_runs)
    assert any(run.text == "return" for run in second_line_runs)
    # keyword run should be colored with the theme's keyword color
    def_run = next(run for run in first_line_runs if run.text == "def")
    assert str(def_run.font.color.rgb) == "0B3D91"


def test_table_renders_grid_and_header(primitives_deck, tmp_path: Path):
    source, asset_root = primitives_deck
    output = tmp_path / "table.pptx"
    build(source, output, asset_root=asset_root)
    prs = Presentation(output)
    slide = prs.slides[1]
    table_shape = _shape(slide, "DSH_TABLE")
    assert table_shape.has_table
    table = table_shape.table
    assert len(table.rows) == 3  # 1 header + 2 body
    assert len(table.columns) == 2
    assert table.cell(0, 0).text == "Method"
    assert table.cell(1, 0).text == "Baseline"
    assert table.cell(2, 1).text == "0.9"
    assert _shape(slide, "DSH_CAPTION").text == "Table 1"


def test_equation_renders_omml_when_pandoc_available(primitives_deck, tmp_path: Path):
    source, asset_root = primitives_deck
    output = tmp_path / "eq.pptx"
    build(source, output, asset_root=asset_root)
    # With pandoc installed, the equation renders as a native <m:oMath> element.
    MATH_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
    omath = _equation_omath(output, 3)
    # Structure: a fraction (m:f) and an n-ary sum (m:nary) with a subscript (m:sSub).
    tags = {etree.QName(node).localname for node in omath.iter() if etree.QName(node).namespace == MATH_NS}
    assert "f" in tags       # \frac{1}{N}
    assert "nary" in tags    # \sum
    assert "sSub" in tags    # y_i


@pytest.mark.skipif(__import__("shutil").which("pandoc") is None, reason="pandoc not installed")
def test_equation_math_fonts_use_drawingml_run_properties(primitives_deck, tmp_path: Path):
    source, asset_root = primitives_deck
    output = tmp_path / "eq-fonts.pptx"
    build(source, output, asset_root=asset_root)
    math_ns = "http://schemas.openxmlformats.org/officeDocument/2006/math"
    drawing_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    omath = _equation_omath(output, 3)
    runs = list(omath.iter("{%s}r" % math_ns))
    assert runs
    for run in runs:
        math_rpr = run.find("{%s}rPr" % math_ns)
        assert math_rpr is None or not any(child.tag.startswith("{%s}" % drawing_ns) for child in math_rpr)
        drawing_rpr = run.find("{%s}rPr" % drawing_ns)
        assert drawing_rpr is not None
        assert drawing_rpr.find("{%s}latin" % drawing_ns).get("typeface") == "Helvetica"
        assert drawing_rpr.find("{%s}ea" % drawing_ns).get("typeface") == "苹方-简"
        assert drawing_rpr.find("{%s}cs" % drawing_ns).get("typeface") == "苹方-简"


@pytest.mark.skipif(__import__("shutil").which("pandoc") is None, reason="pandoc not installed")
def test_equation_zip_xml_is_renderable_and_namespace_normalization_is_idempotent(primitives_deck, tmp_path: Path):
    source, asset_root = primitives_deck
    output = tmp_path / "eq-xml.pptx"
    build(source, output, asset_root=asset_root)

    with zipfile.ZipFile(output) as archive:
        slide_xml = archive.read("ppt/slides/slide3.xml")
    root = etree.fromstring(slide_xml)
    math_ns = "http://schemas.openxmlformats.org/officeDocument/2006/math"
    a14_ns = "http://schemas.microsoft.com/office/drawing/2010/main"
    mc_ns = "http://schemas.openxmlformats.org/markup-compatibility/2006"
    drawing_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    assert root.nsmap["m"] == math_ns
    assert root.nsmap["a14"] == a14_ns
    assert root.nsmap["mc"] == mc_ns
    omath = root.find(".//{%s}oMath" % math_ns)
    assert omath is not None
    for math_rpr in omath.iter("{%s}rPr" % math_ns):
        assert not any(etree.QName(child).namespace == drawing_ns for child in math_rpr)
    assert normalize_math_namespaces_in_pptx(output) is False


@pytest.mark.skipif(__import__("shutil").which("pandoc") is None, reason="pandoc not installed")
def test_equation_uses_powerpoint_math_extension_wrapper(primitives_deck, tmp_path: Path):
    source, asset_root = primitives_deck
    output = tmp_path / "eq-wrapper.pptx"
    build(source, output, asset_root=asset_root)
    with zipfile.ZipFile(output) as archive:
        root = etree.fromstring(archive.read("ppt/slides/slide3.xml"))

    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    mc_ns = "http://schemas.openxmlformats.org/markup-compatibility/2006"
    a14_ns = "http://schemas.microsoft.com/office/drawing/2010/main"
    math_ns = "http://schemas.openxmlformats.org/officeDocument/2006/math"
    choices = root.findall(".//{%s}Choice" % mc_ns)
    assert choices, "PowerPoint equations must be inside mc:AlternateContent"
    alternate = choices[0].getparent()
    fallback = alternate.find("{%s}Fallback" % mc_ns)
    assert fallback is not None
    math_extensions = root.findall(".//{%s}m" % a14_ns)
    assert math_extensions, "expected the a14:m PowerPoint math extension"
    assert math_extensions[0].find("{%s}oMathPara" % math_ns) is not None
    assert fallback.find(".//{%s}m" % a14_ns) is None
    assert root.find(".//{%s}sp/{%s}txBody/{%s}p/{%s}oMathPara" % (p_ns, p_ns, "http://schemas.openxmlformats.org/drawingml/2006/main", math_ns)) is None


@pytest.mark.skipif(__import__("shutil").which("pandoc") is None, reason="pandoc not installed")
def test_equation_matrix_renders_as_omml(tmp_path: Path):
    source = tmp_path / "matrix.yml"
    source.write_text(
        r"""
metadata: {title: Matrix}
sections:
  - id: m
    title: M
    short_title: M
    slides:
      - kind: equation
        title: Matrix
        equation: |
          \begin{pmatrix} a & b \\ c & d \end{pmatrix}
""".strip(),
        encoding="utf-8",
    )
    output = tmp_path / "matrix.pptx"
    build(source, output, asset_root=tmp_path)
    MATH_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
    omath = _equation_omath(output, 1)
    tags = {etree.QName(node).localname for node in omath.iter() if etree.QName(node).namespace == MATH_NS}
    assert "d" in tags and "m" in tags  # delimiter (pmatrix) + matrix


def test_blocks_render_panel_and_title_for_each_style(primitives_deck, tmp_path: Path):
    source, asset_root = primitives_deck
    output = tmp_path / "blocks.pptx"
    build(source, output, asset_root=asset_root)
    prs = Presentation(output)
    for index, expected_title in enumerate(["Note", "Warning", "Demo"]):
        slide = prs.slides[3 + index]
        _shape(slide, "DSH_BLOCK_PANEL")
        title = _shape(slide, "DSH_BLOCK_TITLE")
        assert title.text == expected_title


def test_pptx_zip_valid_for_new_kinds(primitives_deck, tmp_path: Path):
    source, asset_root = primitives_deck
    output = tmp_path / "zip.pptx"
    build(source, output, asset_root=asset_root)
    assert zipfile.is_zipfile(output)
    with zipfile.ZipFile(output) as archive:
        names = set(archive.namelist())
    assert "[Content_Types].xml" in names
    assert "ppt/slides/slide1.xml" in names


@pytest.mark.skipif(
    __import__("shutil").which("xelatex") is None or __import__("shutil").which("pdftocairo") is None,
    reason="TikZ toolchain not installed",
)
def test_figure_slide_embeds_compiled_png(tmp_path: Path):
    """A figure slide compiles its TikZ source at build time and embeds a PNG.

    The slide needs no pre-existing image file — the figure spec is the source
    of truth, so the deck and its figures can never drift apart.
    """
    source = tmp_path / "deck.yml"
    source.write_text(
        r"""
metadata: {title: Fig}
sections:
  - id: a
    title: A
    short_title: A
    slides:
      - kind: title
        title: Fig
      - kind: figure
        title: MLP
        caption: Fig 1
        figure:
          engine: tikz
          source: |
            \begin{tikzpicture}
            \node[draw] (a) {Input}; \node[draw, right=of a] (b) {Out};
            \draw[-Latex] (a) -- (b);
            \end{tikzpicture}
          preamble: ["\\usetikzlibrary{positioning,arrows.meta}"]
""".strip(),
        encoding="utf-8",
    )
    output = tmp_path / "fig.pptx"
    report = build(source, output, asset_root=tmp_path)
    assert report.slide_count == 2
    assert output.is_file()

    prs = Presentation(output)
    figure_slide = prs.slides[1]
    pictures = [s for s in figure_slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pictures, "figure slide must embed a compiled PNG"
    # The compiled figure is named DSH_FIGURE (other images use DSH_IMAGE).
    assert any(p.name == "DSH_FIGURE" for p in pictures)
