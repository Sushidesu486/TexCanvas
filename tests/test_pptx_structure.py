import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from texcanvas import build


def _shape(slide, name):
    return next(shape for shape in slide.shapes if shape.name == name)


def test_generated_pptx_has_editable_named_shapes(deck_files: tuple[Path, Path], tmp_path: Path):
    source, asset_root = deck_files
    output = tmp_path / "structure.pptx"
    build(source, output, asset_root=asset_root)
    prs = Presentation(output)

    assert _shape(prs.slides[1], "DSH_TITLE").text == "Question"
    active_nav = _shape(prs.slides[1], "DSH_NAV_background")
    inactive_nav = _shape(prs.slides[1], "DSH_NAV_methods")
    assert str(active_nav.fill.fore_color.rgb) == "16324F"
    assert str(inactive_nav.fill.fore_color.rgb) == "D9E3EC"
    assert _shape(prs.slides[1], "DSH_FOOTER_PAGE").text == "02 / 06"
    assert any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.name == "DSH_IMAGE" for shape in prs.slides[3].shapes)
    assert "Doe. Paper." in _shape(prs.slides[5], "DSH_REFERENCES").text


def test_pptx_zip_has_required_ooxml_parts(deck_files: tuple[Path, Path], tmp_path: Path):
    source, asset_root = deck_files
    output = tmp_path / "zip-smoke.pptx"
    build(source, output, asset_root=asset_root)
    assert zipfile.is_zipfile(output)
    with zipfile.ZipFile(output) as archive:
        names = set(archive.namelist())
    assert "[Content_Types].xml" in names
    assert "ppt/presentation.xml" in names
    assert "ppt/slides/slide1.xml" in names

