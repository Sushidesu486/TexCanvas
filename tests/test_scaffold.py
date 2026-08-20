from __future__ import annotations

import os
import stat
import subprocess
import sys
from pathlib import Path

import pytest
from pptx import Presentation

from texcanvas import build, init_project
from texcanvas.cli import main
from texcanvas.scaffold import bundled_template_path


def _venv_bin() -> str | None:
    # The editable-install venv holds the `texcanvas` console script used by build.sh.
    candidate = Path(sys.executable).parent / "texcanvas"
    return str(candidate) if candidate.exists() else None


def test_init_project_creates_scaffold(tmp_path: Path):
    project = init_project(tmp_path, "my-talk")
    assert project == (tmp_path / "my-talk").resolve()
    # Full file tree.
    assert (project / "AGENTS.md").is_file()
    assert (project / "deck.yml").is_file()
    assert (project / "build.sh").is_file()
    assert (project / "assets").is_dir()
    assert (project / "output").is_dir()
    assert (project / "output" / ".gitkeep").is_file()
    assert (project / ".gitignore").is_file()
    # build.sh is executable.
    mode = (project / "build.sh").stat().st_mode
    assert mode & stat.S_IXUSR, "build.sh should be executable"
    # deck.yml builds green out of the box (no template arg — uses bundled default).
    report = build(project / "deck.yml", project / "output" / "deck.pptx", asset_root=project)
    assert report.slide_count == 3
    assert report.section_count == 1
    assert report.warnings == ()
    prs = Presentation(project / "output" / "deck.pptx")
    assert len(prs.slides) == 3


def test_init_refuses_existing_destination(tmp_path: Path):
    (tmp_path / "exists").mkdir()
    from texcanvas.errors import InputError

    with pytest.raises(InputError, match="already exists"):
        init_project(tmp_path, "exists")


def test_init_deck_uses_bundled_template_background(tmp_path: Path):
    # No -t passed: the scaffolded deck.yml must still get the F7F9FC background via
    # the bundled template's master, without painting a DSH_BACKGROUND shape.
    project = init_project(tmp_path, "t")
    build(project / "deck.yml", project / "o.pptx", asset_root=project)
    prs = Presentation(project / "o.pptx")
    # slide[1] is a section_divider (no DSH_TITLE); slide[2] is content (has DSH_TITLE).
    names = {sh.name for sh in prs.slides[1].shapes}
    assert "DSH_BACKGROUND" not in names
    content_names = {sh.name for sh in prs.slides[2].shapes}
    assert "DSH_TITLE" in content_names  # content slide chrome still present
    assert "DSH_BACKGROUND" not in content_names


def test_cli_init_subcommand(tmp_path: Path, capsys):
    status = main(["init", "deck-proj", "-d", str(tmp_path)])
    out = capsys.readouterr().out
    assert status == 0
    assert (tmp_path / "deck-proj" / "deck.yml").is_file()
    assert (tmp_path / "deck-proj" / "build.sh").is_file()
    assert "Created" in out
    assert "bash build.sh" in out


def test_bundled_template_path_resolves():
    path = bundled_template_path()
    assert path.is_file()
    assert path.suffix == ".pptx"


@pytest.mark.skipif(_venv_bin() is None, reason="texcanvas console script not on disk")
def test_build_sh_wrapper_runs_end_to_end(tmp_path: Path):
    project = init_project(tmp_path, "e2e")
    env = os.environ.copy()
    venv_bin = Path(sys.executable).parent
    env["PATH"] = f"{venv_bin}:{env.get('PATH', '')}"
    result = subprocess.run(["bash", "build.sh"], cwd=project, env=env, capture_output=True, text=True)
    assert result.returncode == 0, result.stderr
    assert (project / "output" / "deck.pptx").is_file()
    assert "Slides: 3" in result.stdout
