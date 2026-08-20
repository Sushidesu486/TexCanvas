import hashlib
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from beamer_pptx import build


def digest(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def test_template_size_is_preserved_and_source_is_unchanged(deck_files: tuple[Path, Path], tmp_path: Path):
    source, asset_root = deck_files
    template = tmp_path / "template.pptx"
    template_prs = Presentation()
    template_prs.slide_width = Inches(12)
    template_prs.slide_height = Inches(6.75)
    sample = template_prs.slides.add_slide(template_prs.slide_layouts[6])
    sample.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1)).text = "remove me"
    template_prs.save(template)
    before = digest(template)

    output = tmp_path / "templated.pptx"
    build(source, output, template=template, asset_root=asset_root)
    after = digest(template)
    result = Presentation(output)

    assert before == after
    assert result.slide_width == Inches(12)
    assert result.slide_height == Inches(6.75)
    assert len(result.slides) == 6
    assert all("remove me" not in shape.text for slide in result.slides for shape in slide.shapes if hasattr(shape, "text"))

