from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn

from texcanvas import build


def _shape(slide, name):
    return next(shape for shape in slide.shapes if shape.name == name)


def _run_fonts(run):
    rPr = run.font._rPr
    latin = rPr.find(qn("a:latin")) if rPr is not None else None
    ea = rPr.find(qn("a:ea")) if rPr is not None else None
    cs = rPr.find(qn("a:cs")) if rPr is not None else None
    return {
        "latin": latin.get("typeface") if latin is not None else None,
        "ea": ea.get("typeface") if ea is not None else None,
        "cs": cs.get("typeface") if cs is not None else None,
    }


_TITLE_DECK = """
metadata:
  title: 科研训练汇报
  subtitle: A Beamer-style editable presentation
  author: 张三
  institute: Zhejiang University
  date: 2026-03-22
sections:
  - id: cover
    title: Cover
    short_title: Cover
    slides:
      - kind: title
        title: 科研训练汇报
        subtitle: A Beamer-style editable presentation
"""


def test_title_slide_renders_band_title_subtitle_and_byline(tmp_path: Path):
    source = tmp_path / "deck.yml"
    source.write_text(_TITLE_DECK.strip(), encoding="utf-8")
    output = tmp_path / "title.pptx"
    build(source, output, asset_root=tmp_path)

    prs = Presentation(output)
    assert len(prs.slides) == 1
    slide = prs.slides[0]

    band = _shape(slide, "DSH_TITLE_BAND")
    assert band.auto_shape_type is MSO_AUTO_SHAPE_TYPE.RECTANGLE
    assert str(band.fill.fore_color.rgb) == "16324F"

    cover = _shape(slide, "DSH_TITLE_COVER")
    assert cover.text == "科研训练汇报"

    subtitle = _shape(slide, "DSH_TITLE_SUBTITLE")
    assert subtitle.text.startswith("A Beamer-style")

    byline = _shape(slide, "DSH_TITLE_BYLINE")
    # author · institute · date all joined.
    assert "张三" in byline.text
    assert "Zhejiang University" in byline.text
    assert "2026-03-22" in byline.text


def test_title_slide_suppresses_chrome(tmp_path: Path):
    source = tmp_path / "deck.yml"
    source.write_text(_TITLE_DECK.strip(), encoding="utf-8")
    output = tmp_path / "chrome.pptx"
    build(source, output, asset_root=tmp_path)

    slide = Presentation(output).slides[0]
    names = {shape.name for shape in slide.shapes}
    assert not any(name.startswith("DSH_NAV_") for name in names)
    assert "DSH_TITLE_RULE" not in names
    assert "DSH_FOOTER_PAGE" not in names
    assert "DSH_FOOTER_LEFT" not in names


def test_title_slide_falls_back_to_metadata(tmp_path: Path):
    # No slide-level title/subtitle → drawn from metadata.
    deck = _TITLE_DECK.replace("        title: 科研训练汇报\n        subtitle: A Beamer-style editable presentation\n", "")
    source = tmp_path / "deck.yml"
    source.write_text(deck.strip(), encoding="utf-8")
    output = tmp_path / "fallback.pptx"
    build(source, output, asset_root=tmp_path)

    cover = _shape(Presentation(output).slides[0], "DSH_TITLE_COVER")
    assert cover.text == "科研训练汇报"


_FONT_DECK = """
metadata: {title: Fonts}
sections:
  - id: a
    title: A
    short_title: A
    slides:
      - kind: content
        title: 中文标题
        body: 中文正文 English body 123
        bullets: [一项, two]
      - kind: table
        title: 表
        table:
          header: [方法, Score]
          rows: [[Baseline, 0.8]]
      - kind: code
        title: 代码
        code: {lang: python, source: "def f():\n    pass\n"}
      - kind: equation
        title: 公式
        equation: "L = -\\frac{1}{N}"
"""


def test_cjk_runs_use_pingfang_latin_runs_use_helvetica(tmp_path: Path):
    source = tmp_path / "deck.yml"
    source.write_text(_FONT_DECK.strip(), encoding="utf-8")
    output = tmp_path / "fonts.pptx"
    build(source, output, asset_root=tmp_path)

    prs = Presentation(output)

    # Content slide: Chinese body run keeps both latin and ea as 苹方-简
    body = _shape(prs.slides[0], "DSH_BODY")
    fonts = _run_fonts(body.text_frame.paragraphs[0].runs[0])
    assert fonts["ea"] == "苹方-简"
    assert fonts["cs"] == "苹方-简"

    # Table header cell run: latin/ea/cs all set
    table = _shape(prs.slides[1], "DSH_TABLE").table
    cell_fonts = _run_fonts(table.cell(0, 0).text_frame.paragraphs[0].runs[0])
    assert cell_fonts["latin"] == "苹方-简"
    assert cell_fonts["ea"] == "苹方-简"
    assert cell_fonts["cs"] == "苹方-简"

    # Code body: Latin font is Menlo, ea stays 苹方-简 (so CJK in comments renders)
    code = _shape(prs.slides[2], "DSH_CODE_BODY")
    code_fonts = _run_fonts(code.text_frame.paragraphs[0].runs[0])
    assert code_fonts["latin"] == "Menlo"
    assert code_fonts["ea"] == "苹方-简"

    # Equation: Latin font is Helvetica. With pandoc available the equation is a
    # native <m:oMath> element; check the math runs carry the deck fonts.
    eq = _shape(prs.slides[3], "DSH_EQUATION")
    MATH_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
    A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    omaths = list(eq.text_frame._txBody.iter("{%s}oMath" % MATH_NS))
    assert omaths, "expected a native OMML equation when pandoc is available"
    run = next(omaths[0].iter("{%s}r" % MATH_NS))
    rPr = run.find("{%s}rPr" % MATH_NS)
    latin = rPr.find("{%s}latin" % A_NS) if rPr is not None else None
    ea = rPr.find("{%s}ea" % A_NS) if rPr is not None else None
    assert latin is not None and latin.get("typeface") == "Helvetica"
    assert ea is not None and ea.get("typeface") == "苹方-简"


_SQUARE_DECK = """
metadata: {title: Squares}
sections:
  - id: a
    title: A
    short_title: A
    slides:
      - kind: block
        title: B
        block: {style: default, title: Note, body: x}
      - kind: code
        title: C
        code: {lang: python, source: "x = 1\n"}
      - kind: equation
        title: E
        equation: "a = b"
      - kind: two_columns
        title: T
        left: {heading: L, body: x}
        right: {heading: R, body: y}
      - kind: conclusion
        title: End
        takeaway: done
"""


def test_panels_use_rectangular_not_rounded_edges(tmp_path: Path):
    source = tmp_path / "deck.yml"
    source.write_text(_SQUARE_DECK.strip(), encoding="utf-8")
    output = tmp_path / "square.pptx"
    build(source, output, asset_root=tmp_path)

    prs = Presentation(output)
    expected = [
        ("DSH_BLOCK_PANEL", 0),
        ("DSH_CODE_PANEL", 1),
        ("DSH_EQUATION_PANEL", 2),
        ("DSH_LEFT_PANEL", 3),
        ("DSH_RIGHT_PANEL", 3),
        ("DSH_TAKEAWAY_PANEL", 4),
    ]
    for name, slide_index in expected:
        shape = _shape(prs.slides[slide_index], name)
        assert shape.auto_shape_type is MSO_AUTO_SHAPE_TYPE.RECTANGLE, f"{name} should be a plain rectangle"
